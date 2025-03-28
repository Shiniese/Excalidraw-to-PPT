import os
import json
import base64
from PIL import Image
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches


def decode_image(data_url):
    # 去掉 "data:image/png;base64,"
    base64_image_data = data_url.split(',')[1]
    # 解码 base64 数据
    image_data = base64.b64decode(base64_image_data)
    # 打开图片并返回
    return Image.open(BytesIO(image_data))

def decode_and_save_images(filename):
    # 读取.excalidraw文件内容
    with open(filename, 'r') as f:
        data = json.load(f)
    
    # 假设"dataURL"是在"data.elements"这个数组中的某个字典里
    elements = data.get('files', [])
    
    # 创建与.excalidraw文件同名的文件夹
    folder_name = filename[:-len('.excalidraw')]
    if not os.path.exists(folder_name):
        os.makedirs(folder_name)
    
    # 存储图片序号
    image_index = 1
    image_paths = []

    for element in elements:
        if "dataURL" in elements[element] and elements[element]["dataURL"].startswith("data:image/png;base64,"):
            image = decode_image(elements[element]["dataURL"])
            
            image_path = f"{folder_name}/{folder_name}-{image_index}.png"
            image.save(image_path)
            image_paths.append(image_path)
            
            print(f"Saved image: {image_path}")
            image_index += 1

    return image_paths


def add_images_to_pptx(image_paths, prs, slide_title):
    # 添加一个新幻灯片
    slide_layout = prs.slide_layouts[6]  # 使用空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = slide_title
    
    # 计算图片数量和布局
    num_images = len(image_paths)
    cols = int(num_images ** 0.5) + 1
    rows = (num_images + cols - 1) // cols
    
    # 定义幻灯片边界
    slide_width = Inches(8)
    slide_height = Inches(6)
    
    # 添加图片
    left = top = Inches(0.1)  # 留一点边距
    img_margin = Inches(0.1)  # 图片之间的间距
    img_slot_width = (slide_width - img_margin * (cols - 1)) / cols
    img_slot_height = (slide_height - img_margin * (rows - 1)) / rows

    for idx, image_path in enumerate(image_paths):
        # 获取图片的原始尺寸
        img_width, img_height = Image.open(image_path).size
        
        # 计算缩放比例，保持纵横比
        aspect_ratio = img_width / img_height
        if aspect_ratio > (img_slot_width / img_slot_height):
            resized_width = img_slot_width
            resized_height = resized_width / aspect_ratio
        else:
            resized_height = img_slot_height
            resized_width = resized_height * aspect_ratio
        
        # 确保图片不超过分配的空间
        resized_width = min(resized_width, img_slot_width)
        resized_height = min(resized_height, img_slot_height)
        
        # 计算图片的位置
        left_padding = (img_slot_width - resized_width) / 2
        top_padding = (img_slot_height - resized_height) / 2
        
        # 添加图片到幻灯片
        slide.shapes.add_picture(image_path, left + left_padding, top + top_padding, width=resized_width, height=resized_height)
        left += img_slot_width + img_margin
        
        # 如果达到列数限制，移动到下一行
        if (idx + 1) % cols == 0:
            left = Inches(0.1)
            top += img_slot_height + img_margin

def process_excalidraw_files():
    # 创建一个新的PowerPoint演示文稿
    prs = Presentation()
    
    # 获取当前目录下所有的.excalidraw文件
    excalidraw_files = [file for file in os.listdir('.') if file.endswith('.excalidraw')]

    for excalidraw_file in excalidraw_files:
        image_paths = decode_and_save_images(excalidraw_file)

        if image_paths:
            # 将所有图片添加到一个幻灯片中
            slide_title = os.path.splitext(excalidraw_file)[0]
            add_images_to_pptx(image_paths, prs, slide_title)
    
    # 保存PowerPoint文件
    prs.save('excalidraw_to_ppt.pptx')
    print("PowerPoint 文件已保存为 excalidraw_to_ppt.pptx")

if __name__ == "__main__":
    process_excalidraw_files()

print("All images have been processed.")
